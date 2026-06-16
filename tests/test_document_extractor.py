"""DocumentExtractor tests.

Routing + defensive behavior use injected fake handlers (no libs needed). Real
round-trips for pptx/docx run only when the 'extract' extra is installed.
"""

from __future__ import annotations

import io

import pytest

from teg.integrations.files.document_extractor import DocumentExtractor, _clean


def test_routes_by_extension_case_insensitively_and_cleans() -> None:
    extractor = DocumentExtractor(handlers={".pdf": lambda b: "Hello\n\n\n\nworld  "})
    assert extractor.extract("FILE.PDF", b"x") == "Hello\n\nworld"


def test_unsupported_extensions_return_empty() -> None:
    extractor = DocumentExtractor(handlers={".pdf": lambda b: "x"})
    assert extractor.extract("legacy.ppt", b"x") == ""  # legacy binary skipped
    assert extractor.extract("sheet.xlsx", b"x") == ""
    assert extractor.extract("noext", b"x") == ""


def test_handler_failure_is_swallowed() -> None:
    def boom(_: bytes) -> str:
        raise ValueError("corrupt file")

    extractor = DocumentExtractor(handlers={".pdf": boom})
    assert extractor.extract("a.pdf", b"x") == ""  # defensive: never crashes condense


def test_clean_normalizes_and_collapses_blank_lines() -> None:
    assert _clean("caf\xe9\xa0bar\n\n\n\nend   ") == "caf\xe9 bar\n\nend"


def test_real_docx_roundtrip() -> None:
    docx = pytest.importorskip("docx")
    document = docx.Document()
    document.add_paragraph("Hello from docx")
    buffer = io.BytesIO()
    document.save(buffer)

    text = DocumentExtractor().extract("idea.docx", buffer.getvalue())
    assert "Hello from docx" in text


def test_real_pptx_roundtrip() -> None:
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Hello from pptx"
    buffer = io.BytesIO()
    prs.save(buffer)

    text = DocumentExtractor().extract("idea.pptx", buffer.getvalue())
    assert "Hello from pptx" in text
