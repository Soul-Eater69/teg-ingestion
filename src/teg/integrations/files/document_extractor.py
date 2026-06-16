"""Attachment text extraction by format (PDF / PPTX / DOCX).

Routes by file extension to a targeted, permissively-licensed library:
  - .pdf  -> pypdfium2 (PDFium via ctypes; fast and releases the GIL while parsing)
  - .pptx -> python-pptx
  - .docx -> python-docx
Legacy binary .ppt/.doc (and other formats) yield "" - they are skipped and the
caller's min-text filter drops them. Each handler is defensive: any parse failure
returns "" so one bad file never breaks condense. No OCR (image-only docs -> "").

Libraries are imported lazily inside each handler so this module loads even without
the optional 'extract' extra installed (tests inject fake handlers).
"""

from __future__ import annotations

import io
import re
import threading
import unicodedata
from typing import Callable

# PDFium (ctypes) and python-pptx are NOT safe to run from multiple threads at once - two
# concurrent parses deadlock inside the native code (no error, the run just hangs), which is
# why ingest needed --concurrency 1. Serialize only the native parse with this lock: extraction
# runs in asyncio.to_thread worker threads, so one parse at a time across them avoids the
# deadlock while the LLM / embedding / download I/O stays fully parallel. Ticket concurrency can
# go back up - the bottleneck is the per-ticket LLM call, not the now-serialized (fast) parse.
_NATIVE_PARSE_LOCK = threading.Lock()

_ZERO_WIDTH = re.compile("[" + "".join(map(chr, (0x200B, 0x200C, 0x200D, 0xFEFF))) + "]")
_CONTROL = re.compile(r"[\x00-\x08\x0b\x0c\x0e-\x1f]")
_EXTRA_BLANK_LINES = re.compile(r"\n{3,}")


def _extension(filename: str) -> str:
    name = filename.lower().strip()
    dot = name.rfind(".")
    return name[dot:] if dot != -1 else ""


def _extract_pdf(content: bytes) -> str:
    import pypdfium2 as pdfium

    pdf = pdfium.PdfDocument(content)
    try:
        parts: list[str] = []
        for page in pdf:
            textpage = page.get_textpage()
            parts.append(textpage.get_text_range())
            textpage.close()
            page.close()
        return "\n\n".join(parts)
    finally:
        pdf.close()


def _extract_pptx(content: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    parts.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def _extract_docx(content: bytes) -> str:
    import docx

    document = docx.Document(io.BytesIO(content))
    parts = [p.text for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def _default_handlers() -> dict[str, Callable[[bytes], str]]:
    return {".pdf": _extract_pdf, ".pptx": _extract_pptx, ".docx": _extract_docx}


class DocumentExtractor:
    """AttachmentTextExtractor routing by extension; defensive per file."""

    def __init__(self, handlers: dict[str, Callable[[bytes], str]] | None = None) -> None:
        self._handlers = handlers if handlers is not None else _default_handlers()

    def extract(self, filename: str, content: bytes) -> str:
        handler = self._handlers.get(_extension(filename))
        if handler is None:
            return ""  # unsupported (legacy .ppt/.doc, spreadsheets, images, ...)
        try:
            with _NATIVE_PARSE_LOCK:  # one native parse at a time (PDFium/pptx aren't thread-safe)
                text = handler(content)
            return _clean(text)  # cleaning is pure Python - no need to hold the lock
        except Exception:
            return ""  # defensive: a bad/corrupt file contributes nothing, never crashes


def _clean(text: str) -> str:
    text = unicodedata.normalize("NFKC", str(text or ""))
    text = _ZERO_WIDTH.sub("", text)
    text = _CONTROL.sub(" ", text)
    text = text.replace("\xa0", " ")
    text = "\n".join(line.rstrip() for line in text.splitlines())
    return _EXTRA_BLANK_LINES.sub("\n\n", text).strip()


def build_attachment_extractor() -> DocumentExtractor:
    return DocumentExtractor()
